#!/usr/bin/env python3
"""Generate Ndotonic pitch deck as a cinematic PPTX."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "docs" / "pitch-deck"

# Brand palette
BG = RGBColor(0x06, 0x08, 0x07)
BG_DEEP = RGBColor(0x03, 0x04, 0x04)
TEXT = RGBColor(0xE8, 0xEC, 0xE9)
MUTED = RGBColor(0x9A, 0xA8, 0xA0)
GOLD = RGBColor(0xE0, 0xBC, 0x3A)
GOLD_SOFT = RGBColor(0xC9, 0xA2, 0x27)
BRONZE = RGBColor(0xA8, 0x7A, 0x18)
CARD = RGBColor(0x0C, 0x0F, 0x0E)
CARD_EDGE = RGBColor(0x2A, 0x3D, 0x32)
GREEN_GLOW = RGBColor(0x14, 0x22, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.72)
CONTENT_W = Inches(11.9)


class Deck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self.num = 0

    def slide(self):
        self.num += 1
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._bg(slide)
        self._mesh(slide)
        self._accent_bar(slide)
        self._footer(slide, self.num)
        return slide

    def _bg(self, slide, deep=False):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = BG_DEEP if deep else BG

    def _mesh(self, slide):
        for i in range(8):
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(1.2 + i * 1.55), Inches(0), Inches(0.015), SLIDE_H,
            )
            line.fill.solid()
            line.fill.fore_color.rgb = CARD_EDGE
            line.line.fill.background()
        orb = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(10.2), Inches(-1.4), Inches(4.5), Inches(4.5),
        )
        orb.fill.solid()
        orb.fill.fore_color.rgb = GREEN_GLOW
        orb.line.fill.background()
        orb2 = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(-1.8), Inches(5.2), Inches(3.2), Inches(3.2),
        )
        orb2.fill.solid()
        orb2.fill.fore_color.rgb = RGBColor(0x18, 0x14, 0x08)
        orb2.line.fill.background()

    def _accent_bar(self, slide, height=Inches(7.5)):
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.09), height)
        bar.fill.solid()
        bar.fill.fore_color.rgb = GOLD_SOFT
        bar.line.fill.background()
        tick = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.09), Inches(0.62), Inches(0.55), Inches(0.045),
        )
        tick.fill.solid()
        tick.fill.fore_color.rgb = GOLD
        tick.line.fill.background()

    def _footer(self, slide, num):
        rule = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, MARGIN, Inches(6.95), CONTENT_W, Inches(0.012),
        )
        rule.fill.solid()
        rule.fill.fore_color.rgb = CARD_EDGE
        rule.line.fill.background()
        _text(slide, MARGIN, Inches(7.05), Inches(4), Inches(0.3),
              "NDOTONIC", size=9, bold=True, color=BRONZE, tracking=0.18)
        _text(slide, Inches(5.5), Inches(7.05), Inches(2.3), Inches(0.3),
              "glc.cool", size=9, color=MUTED, align=PP_ALIGN.CENTER)
        _text(slide, Inches(11.3), Inches(7.05), Inches(1.3), Inches(0.3),
              f"{num:02d}", size=9, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)

    def _card(self, slide, left, top, width, height, *, fill=CARD, edge=CARD_EDGE, radius=True):
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        card = slide.shapes.add_shape(shape_type, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = fill
        card.line.color.rgb = edge
        card.line.width = Pt(1.25)
        return card

    def _kicker(self, slide, text):
        _text(slide, MARGIN, Inches(0.52), Inches(5), Inches(0.35),
              text.upper(), size=10, bold=True, color=GOLD_SOFT, tracking=0.22)

    def _title(self, slide, text, *, top=Inches(0.88), size=34, width=CONTENT_W):
        _text(slide, MARGIN, top, width, Inches(1.35), text, size=size, bold=True)

    def _subtitle(self, slide, text, top=Inches(2.05), size=16):
        _text(slide, MARGIN, top, CONTENT_W, Inches(1.2), text, size=size, color=MUTED)


def _text(slide, left, top, width, height, text, *, size=18, bold=False,
          color=TEXT, align=PP_ALIGN.LEFT, tracking=0.0, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    if tracking:
        p.font.spacing = Pt(tracking * size)
    return tf


def _bullets(slide, items, left, top, width, height, *, size=14, color=MUTED, gap=10, prefix="◆ "):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{prefix}{item}" if prefix else item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(gap)
    return tf


def _style_cell(cell, text, *, size=13, bold=False, color=TEXT, align=PP_ALIGN.LEFT, fill=None):
    cell.text = text
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Pt(12)
    cell.margin_right = Pt(12)
    cell.margin_top = Pt(8)
    cell.margin_bottom = Pt(8)
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    tf = cell.text_frame
    tf.word_wrap = True
    for p in tf.paragraphs:
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = align


def slide_title(deck):
    slide = deck.slide()
    deck._bg(slide, deep=True)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.55), SLIDE_W, Inches(1.95))
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(0x10, 0x12, 0x10)
    band.line.fill.background()
    glow = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.8), Inches(0.8), Inches(3.8), Inches(3.8))
    glow.fill.solid()
    glow.fill.fore_color.rgb = RGBColor(0x20, 0x18, 0x06)
    glow.line.fill.background()

    _text(slide, Inches(0.7), Inches(1.35), Inches(12), Inches(0.5),
          "NDOTONIC", size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER, tracking=0.35)
    _text(slide, Inches(0.7), Inches(2.05), Inches(12), Inches(2.2),
          "Turning stories into\npublished books", size=48, bold=True, align=PP_ALIGN.CENTER)
    _text(slide, Inches(1.5), Inches(4.35), Inches(10.3), Inches(0.6),
          "Ndotonic builds the AI native author platform.", size=19, color=MUTED, align=PP_ALIGN.CENTER)
    _text(slide, Inches(0.7), Inches(6.15), Inches(12), Inches(0.4),
          "INVESTOR & PARTNER DECK", size=11, bold=True, color=GOLD_SOFT, align=PP_ALIGN.CENTER, tracking=0.2)


def slide_problem(deck):
    slide = deck.slide()
    deck._kicker(slide, "The problem")
    _text(slide, MARGIN, Inches(0.88), CONTENT_W, Inches(1.1),
          "Independent authors are forced to\nstitch together a dozen tools",
          size=28, bold=True)

    deck._card(slide, MARGIN, Inches(2.28), CONTENT_W, Inches(4.35), edge=GOLD_SOFT)
    rows = [
        ("Area", "Pain"),
        ("Discovery & funding", "Campaigns on one platform, pitches on another → Funding campaigns"),
        ("Writing & editing", "Manuscript tools, AI, and human editors if necessary → Ink Studio"),
        ("Publishing", "ISBN, ebook, print, audiobook → Ink Studio"),
        ("Promotion", "Reaching your audience without owned broadcast → GLC Media (radio/TV)"),
        ("Sales", "Digital, audio, and print sales in separate channels → Marketplace"),
    ]
    table = slide.shapes.add_table(
        len(rows), 2, Inches(0.95), Inches(2.48), Inches(11.45), Inches(4.0),
    ).table
    table.columns[0].width = Inches(2.85)
    table.columns[1].width = Inches(8.6)
    for r, row in enumerate(rows):
        table.rows[r].height = Inches(0.48 if r == 0 else 0.7)
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            if r == 0:
                _style_cell(cell, val, size=12, bold=True, color=GOLD, fill=RGBColor(0x14, 0x17, 0x15))
            elif c == 0:
                _style_cell(cell, val, size=12, bold=True, color=GOLD, fill=CARD)
            else:
                _style_cell(cell, val, size=12, color=TEXT, fill=CARD)


def slide_solution(deck):
    slide = deck.slide()
    deck._kicker(slide, "Our solution")
    deck._title(slide, "Ndotonic is the full author stack")
    deck._subtitle(slide,
                   "One company from first pitch to paying reader, with GLC Media built in to promote every story and book.")

    steps = [
        ("Pitch", "Campaign"),
        ("Fund", "Patrons"),
        ("Write", "Ink Studio"),
        ("Edit", "AI craft"),
        ("Publish", "All formats"),
        ("GLC Media", "Radio & TV"),
        ("Monetize", "Marketplace"),
    ]
    x0 = Inches(0.72)
    w = Inches(1.48)
    h = Inches(1.55)
    gap = Inches(0.18)
    y = Inches(3.55)
    for i, (label, sub) in enumerate(steps):
        left = x0 + i * (w + gap)
        deck._card(slide, left, y, w, h, fill=RGBColor(0x0E, 0x12, 0x10))
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, y, w, Inches(0.06))
        accent.fill.solid()
        accent.fill.fore_color.rgb = GOLD if i in (0, 3, 6) else GOLD_SOFT
        accent.line.fill.background()
        _text(slide, left + Inches(0.12), y + Inches(0.22), w - Inches(0.24), Inches(0.5),
              label, size=13, bold=True, color=GOLD)
        _text(slide, left + Inches(0.12), y + Inches(0.72), w - Inches(0.24), Inches(0.55),
              sub, size=10, color=MUTED)
        if i < len(steps) - 1:
            arrow = _text(slide, left + w, y + Inches(0.55), gap, Inches(0.4),
                          "→", size=14, bold=True, color=GOLD_SOFT, align=PP_ALIGN.CENTER)


def slide_pillars(deck):
    slide = deck.slide()
    deck._kicker(slide, "Product")
    deck._title(slide, "Why Ndotonic: four pillars", size=32)
    cards = [
        ("01", "Get Discovered", "Upload a pitch and launch patron book funding campaigns.", GOLD),
        ("02", "Promote Your Story", "GLC Media radio & TV stream promoting stories and books.", GOLD_SOFT),
        ("03", "Self publish", "Ink Studio, AI editing toolkit, ISBN, ebook · print · audiobook.", GOLD),
        ("04", "Monetize your work", "Marketplace digital, audio, and print worldwide.", GOLD_SOFT),
    ]
    positions = [
        (MARGIN, Inches(2.15), Inches(5.75), Inches(2.15)),
        (Inches(6.85), Inches(2.15), Inches(5.75), Inches(2.15)),
        (MARGIN, Inches(4.55), Inches(5.75), Inches(2.15)),
        (Inches(6.85), Inches(4.55), Inches(5.75), Inches(2.15)),
    ]
    for (num, title, body, accent), (left, top, w, h) in zip(cards, positions):
        deck._card(slide, left, top, w, h, fill=RGBColor(0x0A, 0x0D, 0x0B))
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()
        _text(slide, left + Inches(0.28), top + Inches(0.18), Inches(1), Inches(0.35),
              num, size=22, bold=True, color=accent)
        _text(slide, left + Inches(0.28), top + Inches(0.62), w - Inches(0.45), Inches(0.45),
              title, size=16, bold=True)
        _text(slide, left + Inches(0.28), top + Inches(1.15), w - Inches(0.45), Inches(0.85),
              body, size=12, color=MUTED)


def slide_ai(deck):
    slide = deck.slide()
    deck._kicker(slide, "Differentiator")
    deck._title(slide, "AI native author editing", size=32)
    deck._subtitle(slide, "Purpose built inside Ink Studio, not generic chat.", top=Inches(1.95))
    cols = [
        ("Copy editing", "Grammar & punctuation\nSpelling\nLinguistic errors", "01"),
        ("Craft review", "Plot continuity\nPacing & tension\nNarrative style", "02"),
        ("Author workflow", "Chapter editor + versions\nCollaborators & beta readers\nPublishing pipeline", "03"),
    ]
    cw = Inches(3.72)
    for i, (title, body, num) in enumerate(cols):
        left = MARGIN + i * (cw + Inches(0.37))
        top = Inches(2.85)
        deck._card(slide, left, top, cw, Inches(3.35), fill=RGBColor(0x0A, 0x0E, 0x0B))
        _text(slide, left + Inches(0.25), top + Inches(0.2), Inches(0.8), Inches(0.45),
              num, size=26, bold=True, color=GOLD_SOFT)
        _text(slide, left + Inches(0.25), top + Inches(0.75), cw - Inches(0.4), Inches(0.5),
              title, size=15, bold=True, color=GOLD)
        _text(slide, left + Inches(0.25), top + Inches(1.35), cw - Inches(0.4), Inches(1.8),
              body, size=12, color=TEXT)


def slide_market(deck):
    slide = deck.slide()
    deck._kicker(slide, "Market")
    deck._title(slide, "A massive, fragmented creator economy", size=30)
    stats = [
        ("$1.8B+", "US self publishing market\n(~8% CAGR)"),
        ("4M+", "New self published titles\nannually worldwide"),
        ("AI shift", "Authors expect editing tools\ninside their publisher"),
    ]
    sw = Inches(3.72)
    for i, (big, small) in enumerate(stats):
        left = MARGIN + i * (sw + Inches(0.37))
        top = Inches(2.05)
        deck._card(slide, left, top, sw, Inches(2.35), fill=RGBColor(0x0B, 0x0F, 0x0C), edge=GOLD_SOFT)
        _text(slide, left + Inches(0.25), top + Inches(0.35), sw - Inches(0.4), Inches(0.75),
              big, size=34, bold=True, color=GOLD)
        _text(slide, left + Inches(0.25), top + Inches(1.25), sw - Inches(0.4), Inches(0.9),
              small, size=13, color=MUTED)
    deck._card(slide, MARGIN, Inches(4.75), CONTENT_W, Inches(1.35))
    _text(slide, Inches(1.05), Inches(5.05), Inches(10.8), Inches(0.9),
          "Ndotonic targets authors who want discovery, craft tools, multi format publishing, and GLC Media broadcast promotion without five vendors.",
          size=14, color=TEXT)


def slide_business(deck):
    slide = deck.slide()
    deck._kicker(slide, "Business model")
    deck._title(slide, "Revenue at every stage", size=32)
    streams = [
        ("Platform fees", "Marketplace sales: ebook, audio, print, bundles"),
        ("Campaigns & patronage", "Book funding and early supporter investment"),
        ("Publishing services", "ISBN, print, AI audiobook, covers, listing tools"),
        ("GLC Media", "Owned radio & TV promoting stories and books"),
    ]
    for i, (title, body) in enumerate(streams):
        left = MARGIN if i % 2 == 0 else Inches(6.85)
        top = Inches(2.05) + (i // 2) * Inches(2.05)
        w, h = Inches(5.75), Inches(1.75)
        deck._card(slide, left, top, w, h)
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.22), top + Inches(0.22),
                                      Inches(0.55), Inches(0.55))
        pill.fill.solid()
        pill.fill.fore_color.rgb = GOLD_SOFT
        pill.line.fill.background()
        _text(slide, left + Inches(0.22), top + Inches(0.28), Inches(0.55), Inches(0.45),
              str(i + 1), size=14, bold=True, color=BG, align=PP_ALIGN.CENTER)
        _text(slide, left + Inches(0.95), top + Inches(0.22), w - Inches(1.1), Inches(0.4),
              title, size=14, bold=True, color=GOLD)
        _text(slide, left + Inches(0.95), top + Inches(0.68), w - Inches(1.1), Inches(0.85),
              body, size=12, color=MUTED)


def slide_traction(deck):
    slide = deck.slide()
    deck._kicker(slide, "Traction")
    deck._title(slide, "Live at glc.cool", size=36)
    panels = [
        ("Ndotonic platform", [
            "Ink Studio · Patron campaigns · Marketplace",
            "AI editing (6 modes) · Stripe payouts",
            "Publishing pipeline end to end",
        ], MARGIN),
        ("GLC Media", [
            "Radio & TV stream",
            "Live broadcast & author features",
            "Editorial amplification · Publicity tools",
        ], Inches(6.85)),
    ]
    for title, items, left in panels:
        deck._card(slide, left, Inches(2.05), Inches(5.75), Inches(3.85), edge=GOLD_SOFT)
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(2.05), Inches(5.75), Inches(0.62))
        header.fill.solid()
        header.fill.fore_color.rgb = RGBColor(0x16, 0x14, 0x08)
        header.line.fill.background()
        _text(slide, left + Inches(0.28), Inches(2.18), Inches(5.2), Inches(0.4),
              title, size=15, bold=True, color=GOLD)
        _bullets(slide, items, left + Inches(0.35), Inches(2.95), Inches(5.1), Inches(2.7),
                 size=13, color=TEXT, gap=12, prefix="")


def slide_competition(deck):
    slide = deck.slide()
    deck._kicker(slide, "Competition")
    deck._title(slide, "We combine what others split apart", size=30)
    deck._card(slide, MARGIN, Inches(2.0), CONTENT_W, Inches(4.15), edge=GOLD_SOFT)
    rows = [
        ("Capability", "Ndotonic", "KDP", "Kickstarter", "Storyrocket"),
        ("Patron campaigns", "✓", "·", "✓", "✓"),
        ("Writing studio", "✓", "·", "·", "·"),
        ("AI craft editing", "✓", "·", "·", "·"),
        ("Multi format store", "✓", "✓", "·", "·"),
        ("Radio/TV (GLC Media)", "✓", "·", "·", "·"),
    ]
    table = slide.shapes.add_table(
        len(rows), 5, Inches(0.95), Inches(2.22), Inches(11.45), Inches(3.75),
    ).table
    col_w = [Inches(3.5), Inches(2.0), Inches(1.85), Inches(1.85), Inches(2.25)]
    for i, w in enumerate(col_w):
        table.columns[i].width = w
    highlight = RGBColor(0x18, 0x16, 0x08)
    for r, row in enumerate(rows):
        table.rows[r].height = Inches(0.46 if r == 0 else 0.58)
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            if r == 0:
                _style_cell(cell, val, size=11, bold=True, color=GOLD,
                            fill=RGBColor(0x14, 0x17, 0x15), align=PP_ALIGN.CENTER)
            elif c == 1:
                _style_cell(cell, val, size=12, bold=True, color=GOLD,
                            fill=highlight, align=PP_ALIGN.CENTER)
            elif c == 0:
                _style_cell(cell, val, size=11, bold=True, color=TEXT, fill=CARD)
            else:
                color = GOLD if val == "✓" else MUTED
                _style_cell(cell, val, size=12, color=color, fill=CARD, align=PP_ALIGN.CENTER)


def slide_vision(deck):
    slide = deck.slide()
    deck._kicker(slide, "Vision")
    deck._title(slide, "Premium independent publishing", size=32)
    deck._subtitle(slide,
                   "Ndotonic: author platform for audience, AI craft, and funding.\n"
                   "GLC Media: streams and promotes stories and books on radio & TV.",
                   top=Inches(1.95))
    years = [
        ("Year 1", "Author acquisition · campaign GMV"),
        ("Year 2", "GLC Media audience · audio/print scale"),
        ("Year 3", "Rights & adaptation pipeline"),
    ]
    yw = Inches(3.72)
    for i, (label, detail) in enumerate(years):
        left = MARGIN + i * (yw + Inches(0.37))
        top = Inches(3.35)
        deck._card(slide, left, top, yw, Inches(2.15))
        node = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(1.45), top - Inches(0.28),
                                      Inches(0.82), Inches(0.82))
        node.fill.solid()
        node.fill.fore_color.rgb = GOLD if i == 1 else GOLD_SOFT
        node.line.color.rgb = GOLD
        node.line.width = Pt(1.5)
        _text(slide, left, top + Inches(0.55), yw, Inches(0.45),
              label, size=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        _text(slide, left + Inches(0.25), top + Inches(1.05), yw - Inches(0.5), Inches(0.85),
              detail, size=12, color=MUTED, align=PP_ALIGN.CENTER)
        if i < 2:
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left + yw, top + Inches(0.95), Inches(0.37), Inches(0.04),
            )
            line.fill.solid()
            line.fill.fore_color.rgb = GOLD_SOFT
            line.line.fill.background()


def slide_team(deck):
    slide = deck.slide()
    deck._kicker(slide, "Team")
    deck._title(slide, "Ndotonic & GLC Media", size=32)
    entities = [
        ("Ndotonic", "The company: Ink Studio, campaigns, marketplace, AI tools", MARGIN),
        ("GLC Media", "Radio & TV stream promoting stories and books", Inches(6.85)),
    ]
    for name, desc, left in entities:
        deck._card(slide, left, Inches(2.05), Inches(5.75), Inches(2.35), edge=GOLD_SOFT)
        _text(slide, left + Inches(0.3), Inches(2.35), Inches(5), Inches(0.45),
              name, size=20, bold=True, color=GOLD)
        _text(slide, left + Inches(0.3), Inches(2.95), Inches(5.1), Inches(1.1),
              desc, size=13, color=TEXT)
    deck._card(slide, MARGIN, Inches(4.75), CONTENT_W, Inches(1.45), fill=RGBColor(0x10, 0x12, 0x0E))
    _text(slide, Inches(1.05), Inches(5.0), Inches(10.8), Inches(0.45),
          "Mission: Uncover hidden potential and amplify unheard voices.", size=14, bold=True, color=GOLD)
    _text(slide, Inches(1.05), Inches(5.45), Inches(10.8), Inches(0.45),
          "Hiring: Co founder CTO, AI engineers, QA, security", size=13, color=MUTED)


def slide_ask(deck):
    slide = deck.slide()
    deck._bg(slide, deep=True)
    glow = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.5), Inches(-0.5), Inches(6.5), Inches(6.5))
    glow.fill.solid()
    glow.fill.fore_color.rgb = RGBColor(0x22, 0x1A, 0x06)
    glow.line.fill.background()

    deck._kicker(slide, "The ask")
    _text(slide, Inches(0.7), Inches(1.15), Inches(12), Inches(0.9),
          "Partner with Ndotonic", size=40, bold=True, align=PP_ALIGN.CENTER)

    deck._card(slide, Inches(2.2), Inches(2.35), Inches(8.9), Inches(2.55), edge=GOLD)
    _bullets(slide, [
        "Strategic investment: author platform & GLC Media",
        "Publishing & distribution partnerships",
        "Broadcast amplification via GLC Media",
    ], Inches(2.65), Inches(2.65), Inches(8.0), Inches(2.0), size=15, color=TEXT, gap=14)

    contact = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.4), Inches(5.15), Inches(6.5), Inches(1.15))
    contact.fill.solid()
    contact.fill.fore_color.rgb = RGBColor(0x18, 0x14, 0x06)
    contact.line.color.rgb = GOLD
    contact.line.width = Pt(1.5)
    _text(slide, Inches(3.4), Inches(5.38), Inches(6.5), Inches(0.35),
          "Email: info@ndotonic.com", size=15, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    _text(slide, Inches(3.4), Inches(5.78), Inches(6.5), Inches(0.35),
          "Instagram: @ndotonic_", size=14, color=TEXT, align=PP_ALIGN.CENTER)


def build_pptx():
    deck = Deck()
    slide_title(deck)
    slide_problem(deck)
    slide_solution(deck)
    slide_pillars(deck)
    slide_ai(deck)
    slide_market(deck)
    slide_business(deck)
    slide_traction(deck)
    slide_competition(deck)
    slide_vision(deck)
    slide_team(deck)
    slide_ask(deck)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    out = OUT_DIR / "Ndotonic_Pitch_Deck.pptx"
    deck.prs.save(str(out))
    return out


if __name__ == "__main__":
    path = build_pptx()
    print(f"Wrote {path}")
